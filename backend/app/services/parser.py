import io
from pathlib import Path

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader


def extract_text_from_file(path: Path, filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix == ".txt":
        return path.read_text(encoding="utf-8", errors="replace")
    if suffix == ".pdf":
        reader = PdfReader(str(path))
        parts: list[str] = []
        for page in reader.pages:
            parts.append(page.extract_text() or "")
        return "\n".join(parts)
    if suffix == ".docx":
        doc = DocxDocument(str(path))
        return "\n".join(p.text for p in doc.paragraphs)
    if suffix == ".pptx":
        prs = Presentation(str(path))
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)
        return "\n".join(lines)
    raise ValueError(f"不支持的格式: {suffix}")


def extract_text_from_upload(content: bytes, filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix == ".txt":
        return content.decode("utf-8", errors="replace")
    buf = io.BytesIO(content)
    if suffix == ".pdf":
        reader = PdfReader(buf)
        parts: list[str] = []
        for page in reader.pages:
            parts.append(page.extract_text() or "")
        return "\n".join(parts)
    if suffix == ".docx":
        doc = DocxDocument(buf)
        return "\n".join(p.text for p in doc.paragraphs)
    if suffix == ".pptx":
        prs = Presentation(buf)
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)
        return "\n".join(lines)
    raise ValueError(f"不支持的格式: {suffix}")
